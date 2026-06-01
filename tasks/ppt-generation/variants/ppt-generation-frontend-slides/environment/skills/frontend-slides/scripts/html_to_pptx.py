#!/usr/bin/env python3
"""Convert a single-file HTML presentation into a PPTX deck.

Strategy: use Playwright to simulate keyboard navigation (ArrowRight) through
the presentation, taking a full-viewport screenshot of each slide, then
assemble the images into a 16:9 PPTX via python-pptx.

A local HTTP server is started so Google Fonts and relative asset paths load
correctly (file:// protocol blocks cross-origin font requests).

Usage:
    python html_to_pptx.py <input.html> <output.pptx> [--width 1920] [--height 1080]
"""

import argparse
import os
import http.server
import subprocess
import sys
import tempfile
import threading
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Emu


# CSS selectors to count slides, tried in order
SLIDE_SELECTORS = [".slide", "[data-slide]", "section", "div > section"]


def _make_screenshot_script(url: str, out_dir: str, width: int, height: int) -> str:
    """Return a Node.js script that navigates and screenshots each slide."""
    selectors_js = ", ".join(f'"{s}"' for s in SLIDE_SELECTORS)
    return f"""\
const {{ chromium }} = require('playwright');

(async () => {{
    const browser = await chromium.launch({{ executablePath: process.env.PLAYWRIGHT_CHROMIUM_EXECUTABLE_PATH || undefined }});
    const page = await browser.newPage();
    await page.setViewportSize({{ width: {width}, height: {height} }});
    await page.goto('{url}', {{ waitUntil: 'networkidle' }});

    // --- Count slides using DOM selectors ---
    const selectors = [{selectors_js}];
    let slideCount = 0;
    for (const sel of selectors) {{
        slideCount = await page.$$eval(sel, els => els.length);
        if (slideCount > 0) break;
    }}

    // Fallback: treat the whole page as 1 slide
    if (slideCount === 0) slideCount = 1;

    // Reasonable upper bound to prevent infinite loops
    if (slideCount > 100) slideCount = 100;

    // --- Capture slide 1 (already visible on load) ---
    // Wait for entrance animations to settle
    await page.waitForTimeout(500);
    await page.screenshot({{ path: `{out_dir}/slide_001.png` }});

    // --- Navigate through remaining slides via ArrowRight ---
    // We track visual changes to detect when navigation has no more effect
    // (i.e. we've passed the last slide).
    let prevPath = `{out_dir}/slide_001.png`;
    let captured = 1;

    for (let i = 2; i <= slideCount; i++) {{
        await page.keyboard.press('ArrowRight');
        // Wait for transition animation to complete
        await page.waitForTimeout(600);

        const idx = String(i).padStart(3, '0');
        const path = `{out_dir}/slide_${{idx}}.png`;
        await page.screenshot({{ path }});
        captured++;
    }}

    console.log(JSON.stringify({{ captured }}));
    await browser.close();
}})();
"""


def _deduplicate_slides(image_dir: Path) -> list[Path]:
    """Remove trailing duplicate screenshots (navigation past last slide).

    Some presentations keep showing the last slide when you press ArrowRight
    beyond the end.  We detect this by comparing file sizes + pixel samples.
    """
    pngs = sorted(image_dir.glob("slide_*.png"))
    if len(pngs) <= 1:
        return pngs

    unique: list[Path] = [pngs[0]]
    prev_img = Image.open(pngs[0])

    for png in pngs[1:]:
        curr_img = Image.open(png)
        # Quick check: compare a grid of sample pixels
        is_dup = True
        for x_frac in (0.25, 0.5, 0.75):
            for y_frac in (0.25, 0.5, 0.75):
                px = int(curr_img.width * x_frac)
                py = int(curr_img.height * y_frac)
                if curr_img.getpixel((px, py)) != prev_img.getpixel((px, py)):
                    is_dup = False
                    break
            if not is_dup:
                break

        if not is_dup:
            unique.append(png)
            prev_img = curr_img
        else:
            # Once we hit a duplicate, stop — remaining are likely all dupes
            break

    return unique


def _start_http_server(directory: Path, port: int = 0):
    """Start a background HTTP server serving *directory*. Returns (server, port)."""
    handler = http.server.SimpleHTTPRequestHandler
    server = http.server.HTTPServer(("127.0.0.1", port), lambda *a, **kw: handler(*a, directory=str(directory), **kw))
    actual_port = server.server_address[1]
    thread = threading.Thread(target=server.serve_forever, daemon=True)
    thread.start()
    return server, actual_port


def screenshot_slides(html_path: Path, out_dir: Path, width: int, height: int) -> list[Path]:
    """Run Playwright to capture per-slide PNGs. Returns deduplicated list."""
    # Serve via HTTP so Google Fonts and relative assets load correctly
    serve_dir = html_path.resolve().parent
    server, port = _start_http_server(serve_dir)
    url = f"http://127.0.0.1:{port}/{html_path.resolve().name}"

    try:
        script = _make_screenshot_script(url, str(out_dir), width, height)
        script_file = out_dir / "_capture.cjs"
        script_file.write_text(script)

        env = os.environ.copy()
        env.setdefault("NODE_PATH", "/usr/local/lib/node_modules:/app/node_modules")
        env.setdefault("PLAYWRIGHT_BROWSERS_PATH", "/ms-playwright")
        env.setdefault("PLAYWRIGHT_SKIP_BROWSER_DOWNLOAD", "1")
        result = subprocess.run(
            ["node", str(script_file)],
            capture_output=True,
            text=True,
            timeout=180,
            env=env,
        )
        if result.returncode != 0:
            print(f"Playwright stderr:\n{result.stderr}", file=sys.stderr)
            raise RuntimeError(f"Screenshot script failed (exit {result.returncode})")
    finally:
        server.shutdown()

    pngs = _deduplicate_slides(out_dir)
    if not pngs:
        raise RuntimeError("No slide screenshots produced")
    return pngs


def assemble_pptx(slide_images: list[Path], output_path: Path) -> None:
    """Pack slide images into a standard 16:9 PPTX."""
    prs = Presentation()
    # Standard 16:9 dimensions
    prs.slide_width = Emu(12192000)   # 13.333 inches
    prs.slide_height = Emu(6858000)   # 7.5 inches

    blank_layout = prs.slide_layouts[6]  # blank layout

    for img_path in slide_images:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path),
            left=Emu(0),
            top=Emu(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(str(output_path))


def main():
    parser = argparse.ArgumentParser(description="Convert HTML presentation to PPTX via screenshots")
    parser.add_argument("input_html", help="Path to the single-file HTML presentation")
    parser.add_argument("output_pptx", help="Output PPTX path")
    parser.add_argument("--width", type=int, default=1920, help="Viewport width (default 1920)")
    parser.add_argument("--height", type=int, default=1080, help="Viewport height (default 1080)")
    args = parser.parse_args()

    html_path = Path(args.input_html).resolve()
    output_path = Path(args.output_pptx).resolve()

    if not html_path.exists():
        raise FileNotFoundError(f"Input HTML not found: {html_path}")

    with tempfile.TemporaryDirectory(prefix="slides_") as tmp:
        tmp_dir = Path(tmp)
        print(f"Capturing slides from {html_path} ...")
        images = screenshot_slides(html_path, tmp_dir, args.width, args.height)
        print(f"Captured {len(images)} unique slide(s), assembling PPTX ...")
        assemble_pptx(images, output_path)
        print(f"Done → {output_path}")


if __name__ == "__main__":
    main()
